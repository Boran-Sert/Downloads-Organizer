import fitz
from pathlib import Path
from docx import Document


class FeatureExtractor:
    """
    Dosyalardan makine öğrenimi modelleri için anlamlı metin (özellik) çıkaran modül.
    """

    @staticmethod
    def extract_from_name(file_name: str) -> str:
        """
        Dosya adındaki özel karakterleri temizler ve NLP için kelimelere ayırır.
        Örn: '2026_Yillik_Finans-Raporu_v2.pdf' -> '2026 Yillik Finans Raporu v2 pdf'
        """
        clean_name = file_name.replace("_", " ").replace("-", " ").replace(".", " ")
        return " ".join(clean_name.split())

    @staticmethod
    def extract_from_pdf(file_path: str) -> str:
        """
        PyMuPDF kullanarak PDF dosyasının içeriğini okur ve tek bir metin olarak döndürür.
        """
        ex_text = []
        try:
            with fitz.open(file_path) as doc:
                for page in doc:
                    ex_text.append(page.get_text())
        except Exception:
            return ""

        return "".join(ex_text)

    @staticmethod
    def extract_from_docx(file_path: str) -> str:
        """
        python-docx kullanarak Word dosyasının içeriğini okur ve tek bir metin olarak döndürür.
        """
        ex_text = []
        try:
            doc = Document(file_path)
            for para in doc.paragraphs:
                ex_text.append(para.text)
        except Exception:
            return ""

        return "".join(ex_text)

    @staticmethod
    def extract_from_pptx(file_path: str) -> str:
        """
        Modern PowerPoint (.pptx) sunumlarının içindeki tüm slaytları gezer
        ve içerideki metin kutularından yazıları çıkarır.
        """
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_runs = []

            # Her bir slayta gir
            for slide in prs.slides:
                # Slaytın içindeki her bir şekli (metin kutusu, başlık vb.) kontrol et
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)

            # Çıkarılan tüm metinleri tek bir uzun metin olarak birleştir
            return " ".join(text_runs)
        except Exception as e:
            # Okunamayan veya bozuk slaytlar programı çökertmesin
            print(f"Uyarı: PPTX okuma hatası ({Path(file_path).name}): {e}")
            return ""

    @staticmethod
    def extract_from_image(file_path: str) -> str:
        """
        Görüntü (png, jpg) içindeki metinleri okur (OCR).
        """
        try:
            import pytesseract
            from PIL import Image

            # Tesseract OCR'ın yüklü olduğu mutlak yolu buraya girmek gerekebilir.
            # Örneğin (Varsayılan Windows kurulum yolu):
            pytesseract.pytesseract.tesseract_cmd = (
                r"C:\Program Files\Tesseract-OCR\tesseract.exe"
            )

            # Görüntüyü aç
            img = Image.open(file_path)

            # Metni oku (Türkçe ve İngilizce karakterler için lang='tur+eng')
            text = pytesseract.image_to_string(img, lang="tur+eng")

            return text
        except Exception as e:
            print(f"Uyarı: Resim okuma hatası ({Path(file_path).name}): {e}")
            return ""

    @staticmethod
    def has_face(file_path: str) -> bool:
        """OpenCV ile resimde insan yüzü olup olmadığını kontrol eder."""
        try:
            import cv2

            # OpenCV'nin varsayılan ön yüz tanıma modelini yükle
            face_cascade = cv2.CascadeClassifier(
                cv2.data.haarcascades + "haarcascade_frontalface_default.xml"
            )

            # Resmi oku ve bilgisayarın daha hızlı işlemesi için gri tonlamaya çevir
            img = cv2.imread(file_path)
            if img is None:
                return False
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

            # Resimdeki yüzleri tara
            faces = face_cascade.detectMultiScale(
                gray, scaleFactor=1.1, minNeighbors=5, minSize=(30, 30)
            )

            # Eğer en az 1 yüz bulunduysa True, bulunamadıysa False dön
            return len(faces) > 0
        except Exception as e:
            print(f"Uyarı: Yüz tarama hatası ({Path(file_path).name}): {e}")
            return False
